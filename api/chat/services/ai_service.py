import time
import uuid
import json
import os
import re
import httpx
import anyio
import asyncio
from sqlalchemy import update
from typing import Optional, List, Dict
from ..core.config import GIGACHAT_KEY, GIGACHAT_MODEL, CERT_VERIFY, log, BASE_DIR
from ..services.db_service import AsyncSessionLocal, ChatMessage, User, update_user_balance
from .gigachat_service import handle_gigachat_request, get_gigachat_token
from .yandex_service import handle_yandex_request
from .data_guard import apply_data_guard
from .vector_service import VectorService
from .response_validator import validate_response_with_agent

def sanitize_history(messages: List[Dict], client_config) -> List[Dict]:
    """
    Очищает историю диалога от устаревших фактов:
    - Телефоны, email, адреса
    - Мессенджеры (Telegram, WhatsApp, Max, VK)
    - Соцсети (extra_links)
    - Реквизиты (ИНН, ОГРН, счета, БИК)
    - Сайт
    которые пользователь мог упомянуть в своих сообщениях.
    """
    import copy
    sanitized = copy.deepcopy(messages)
    
    if client_config is None:
        client_config = {}
    
    is_dict = isinstance(client_config, dict)
    contacts = client_config.get('contacts', {}) if is_dict else client_config.raw.get('contacts', {})
    site_url = client_config.get('site_url', '') if is_dict else client_config.site_url
    legal_data = client_config.get('legal_data', {}) if is_dict else client_config.raw.get('legal_data', {})
    legal_cfg = client_config.get('legal', {}) if is_dict else client_config.raw.get('legal', {})
    
    # ═══════════════════════════════════════════
    # БЕЛЫЙ СПИСОК АКТУАЛЬНЫХ ДАННЫХ
    # ═══════════════════════════════════════════
    current_facts = set()
    
    # --- Телефоны ---
    main_phone = contacts.get('phone', '')
    if main_phone:
        current_facts.add(re.sub(r'[^0-9+]', '', main_phone))
    for p in contacts.get('extra_phones', []):
        if isinstance(p, dict) and p.get('phone'):
            current_facts.add(re.sub(r'[^0-9+]', '', p['phone']))
    
    # --- Email ---
    main_email = contacts.get('email', '')
    if main_email:
        current_facts.add(main_email.lower().strip())
    for e in contacts.get('extra_emails', []):
        if isinstance(e, dict) and e.get('email'):
            current_facts.add(e['email'].lower().strip())
    
    # --- Адреса ---
    for a in contacts.get('extra_addresses', []):
        if isinstance(a, dict) and a.get('address'):
            current_facts.add(a['address'].lower().strip())
    
    # --- Мессенджеры: Telegram, WhatsApp, Max, VK ---
    for m_key in ['extra_tg', 'extra_wa', 'extra_max', 'extra_vk']:
        for item in contacts.get(m_key, []):
            if isinstance(item, dict) and item.get('value'):
                val = item['value'].strip()
                current_facts.add(val.lower())
                # Добавляем варианты: с @ и без, t.me/..., wa.me/..., vk.me/...
                clean = val.replace('@', '').replace('https://', '').replace('http://', '')
                current_facts.add(clean.lower())
                if val.startswith('@'):
                    current_facts.add(f"t.me/{val[1:]}".lower())
                if m_key == 'extra_wa':
                    digits = re.sub(r'\D', '', val)
                    if digits:
                        current_facts.add(digits)
                        current_facts.add(f"wa.me/{digits}".lower())
    
    # --- Соцсети (extra_links) ---
    for link in contacts.get('extra_links', []):
        if isinstance(link, dict) and link.get('url'):
            url = link['url'].strip().lower()
            current_facts.add(url)
            # Добавляем без протокола
            clean = url.replace('https://', '').replace('http://', '').rstrip('/')
            current_facts.add(clean)
    
    # --- Сайт ---
    if site_url:
        su = site_url.lower().strip()
        current_facts.add(su)
        current_facts.add(su.replace('https://', '').replace('http://', '').rstrip('/'))
    
    # --- Реквизиты (ИНН, ОГРН, счета, БИК) ---
    legal_type = legal_cfg.get('type', 'ip') if legal_cfg else 'ip'
    type_data = legal_data.get(legal_type, {})
    for key in ['inn', 'ogrn', 'bank_bik', 'bank_account', 'bank_corr', 'pass_seria', 'pass_number']:
        val = type_data.get(key, '')
        if val:
            current_facts.add(str(val).strip())
    
    # --- Название компании ---
    company_name = contacts.get('company_name', '')
    if company_name:
        current_facts.add(company_name.lower().strip())
    
    # ═══════════════════════════════════════════
    # ПАТТЕРНЫ ДЛЯ ПОИСКА ФАКТОВ В ТЕКСТЕ
    # ═══════════════════════════════════════════
    phone_pattern = re.compile(r'(?:\+?[78][\s\-]?\(?\d{3}\)?[\s\-]?\d{3}[\s\-]?\d{2}[\s\-]?\d{2})')
    email_pattern = re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}')
    # Мессенджеры: @username, t.me/..., wa.me/..., vk.me/...
    messenger_pattern = re.compile(
        r'(?:@[a-zA-Z0-9_]{3,32})'  # @username
        r'|(?:t\.me/[a-zA-Z0-9_+]{3,32})'  # t.me/username
        r'|(?:wa\.me/\d{7,15})'  # wa.me/79990000000
        r'|(?:vk\.me/[a-zA-Z0-9_.]{3,32})'  # vk.me/username
        r'|(?:https?://(?:t\.me|wa\.me|vk\.me)/[^\s]+)'  # полные ссылки
    )
    # Соцсети: URL vk.com, youtube.com, instagram.com и др.
    social_url_pattern = re.compile(
        r'https?://(?:vk\.com|youtube\.com|instagram\.com|facebook\.com|t\.me|ok\.ru|dzen\.ru|rutube\.ru)'
        r'/[^\s]+'
    )
    # Реквизиты: ИНН (10 или 12 цифр), ОГРН/ОГРНИП (13 или 15 цифр), БИК (9 цифр), счета (20 цифр)
    requisites_pattern = re.compile(
        r'(?:ИНН[\s:]*\d{10,12})'  # ИНН с меткой
        r'|(?:ОГРН(?:ИП)?[\s:]*\d{13,15})'  # ОГРН/ОГРНИП с меткой
        r'|(?:БИК[\s:]*\d{9})'  # БИК с меткой
        r'|(?:р/с[\s:]*\d{20})'  # расчётный счёт с меткой
        r'|(?:к/с[\s:]*\d{20})'  # корр. счёт с меткой
        r'|(?:\b\d{20}\b)'  # голые 20-значные номера (счета)
        r'|(?:\b\d{9}\b)'  # голые 9-значные (БИК)
    )
    # Адреса: паттерн вида "г. Москва, ул. Ленина, д. 1"
    address_pattern = re.compile(
        r'(?:г\.?\s*[А-ЯЁ][а-яё]+)'  # г. Москва
        r'|(?:ул\.?\s*[А-ЯЁ][а-яё]+)'  # ул. Ленина
        r'|(?:д\.?\s*\d+[а-я]?)'  # д. 1
        r'|(?:пр-т\.?\s*[А-ЯЁ][а-яё]+)'  # пр-т Мира
        r'|(?:пер\.?\s*[А-ЯЁ][а-яё]+)'  # пер. Тихий
        r'|(?:пл\.?\s*[А-ЯЁ][а-яё]+)'  # пл. Ленина
        r'|(?:наб\.?\s*[А-ЯЁ][а-яё]+)'  # наб. реки
        r'|(?:ш\.?\s*[А-ЯЁ][а-яё]+)'  # ш. Энтузиастов
        r'|(?:б-р\.?\s*[А-ЯЁ][а-яё]+)'  # б-р Цветной
    )
    
    # ═══════════════════════════════════════════
    # ОЧИСТКА СООБЩЕНИЙ
    # ═══════════════════════════════════════════
    for i, msg in enumerate(sanitized):
        if msg.get('role') != 'user':
            continue
        if i == len(sanitized) - 1:
            continue  # Последнее сообщение не трогаем
        
        content = msg.get('content', '')
        if not content:
            continue
        
        modified = False
        
        # Телефоны
        for match in phone_pattern.finditer(content):
            phone_digits = re.sub(r'[^0-9+]', '', match.group())
            if phone_digits not in current_facts:
                content = content.replace(match.group(), '[УСТАРЕВШИЙ ТЕЛЕФОН УДАЛЁН]')
                modified = True
        
        # Email
        for match in email_pattern.finditer(content):
            email = match.group().lower().strip()
            if email not in current_facts:
                content = content.replace(match.group(), '[УСТАРЕВШАЯ ПОЧТА УДАЛЕНА]')
                modified = True
        
        # Мессенджеры
        for match in messenger_pattern.finditer(content):
            val = match.group().lower().strip()
            # Проверяем разные варианты
            is_known = False
            for check in [val, val.replace('@', ''), val.replace('https://', '').replace('http://', '')]:
                if check in current_facts:
                    is_known = True
                    break
            if not is_known:
                content = content.replace(match.group(), '[УСТАРЕВШИЙ КОНТАКТ МЕССЕНДЖЕРА УДАЛЁН]')
                modified = True
        
        # Соцсети
        for match in social_url_pattern.finditer(content):
            url = match.group().lower().strip()
            clean = url.replace('https://', '').replace('http://', '').rstrip('/')
            if url not in current_facts and clean not in current_facts:
                content = content.replace(match.group(), '[УСТАРЕВШАЯ ССЫЛКА НА СОЦСЕТЬ УДАЛЕНА]')
                modified = True
        
        # Реквизиты
        for match in requisites_pattern.finditer(content):
            val = match.group().strip()
            digits_only = re.sub(r'[^0-9]', '', val)
            # Проверяем: есть ли эти цифры в белом списке
            is_known = any(digits_only in fact for fact in current_facts if fact)
            if not is_known and digits_only not in current_facts:
                content = content.replace(match.group(), '[УСТАРЕВШИЕ РЕКВИЗИТЫ УДАЛЕНЫ]')
                modified = True
        
        # Адреса (только если в сообщении есть адресные маркеры, но их нет в белом списке)
        for match in address_pattern.finditer(content):
            fragment = match.group().lower().strip()
            # Проверяем, есть ли этот фрагмент в каком-либо из актуальных адресов
            is_known = any(fragment in fact for fact in current_facts if fact)
            if not is_known:
                # Не удаляем отдельные фрагменты (слишком агрессивно),
                # но логируем для отладки
                log.info(f"[SANITIZE] Unrecognized address fragment in msg #{i}: {fragment}")
        
        if modified:
            msg['content'] = content
            log.info(f"[SANITIZE] Cleaned outdated facts from user message #{i}")
    
    return sanitized


def get_bot_tools(bot_settings: dict = None, is_admin: bool = False):
    """Возвращает описание инструментов для ассистента."""
    tools = []
    
    # Функция получения времени доступна всегда, так как не требует интернета
    tools.append({
        "name": "get_datetime",
        "description": "Получить текущую дату и время сервера.",
        "parameters": {"type": "object", "properties": {}}
    })
    return tools

async def ask_ai(messages_raw: List[Dict], client_config=None, model_name=None, stream=False, attachments=None, **kwargs):
    """
    Основная функция запроса к ИИ.
    """
    # Делаем глубокую копию сообщений, чтобы наши системные вставки НЕ попали в БД
    import copy
    messages = copy.deepcopy(messages_raw)

    client_id = kwargs.get('client_id', 'unknown')

    if client_config is None:
        client_config = {}
    
    # Санитизация истории: удаляем устаревшие факты из сообщений пользователя
    messages = sanitize_history(messages, client_config)
    is_dict = isinstance(client_config, dict)
    bot_settings = (client_config.get('bot_settings', {}) if is_dict else client_config.raw.get('bot_settings', {}))
    is_admin = kwargs.get('is_admin', False)

    # 1. СТУПЕНЬ: РЕГЛАМЕНТ И ФАЙЛЫ - ВЫСШИЙ ПРИОРИТЕТ
    # Добавляем "Золотое правило" актуальности в начало промпта
    updated_at_ts = client_config.get('updated_at', 0) if is_dict else client_config.raw.get('updated_at', 0)
    updated_at_str = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(updated_at_ts))
    
    system_instruction = (
        "### [БАЗА ЗНАНИЙ И РЕГЛАМЕНТЫ]\n"
        f"Данные актуальны на: {updated_at_str}\n"
        "Вся информация ниже (контакты, график, реквизиты, база знаний) является приоритетной и единственно верной.\n"
        "Правило приоритета: Если информация в истории диалога противоречит данным из этого реестра — история считается устаревшей. "
        "Всегда используй актуальные данные из разделов ниже.\n\n"
    )
    
    # КОРРЕКЦИЯ ИСТОРИИ: Если конфиг обновился, очищаем старые ответы бота от фактов
    if len(messages) > 1:
        # Проходим по истории и "очищаем" ответы ассистента, которые были даны до обновления конфига
        # Мы не удаляем их совсем, чтобы не ломать логику диалога, но стираем контент
        for msg in messages[:-1]: # Не трогаем последнее сообщение (оно обычно от юзера)
            if msg.get('role') == 'assistant':
                # Если сообщение от бота, мы заменяем его на нейтральную фразу
                # Это заставляет модель брать данные ТОЛЬКО из системного промпта
                msg['content'] = "[ИНФОРМАЦИЯ В ЭТОМ СООБЩЕНИИ УСТАРЕЛА. АКТУАЛЬНЫЕ ДАННЫЕ ВЗЯТЫ ИЗ СИСТЕМНОГО РЕЕСТРА]"
        
        # Добавляем финальное напоминание в последнее сообщение пользователя
        last_msg = messages[-1]
        if last_msg.get('role') == 'user':
            reminder = (
                f"\n\n[Системная заметка: Данные обновлены {updated_at_str}. "
                "Используй только актуальную информацию из системного промпта, игнорируя устаревшие факты из истории диалога.]"
            )
            last_msg['content'] = last_msg.get('content', '') + reminder
    
    file_content = system_instruction
    user_attachments_content = ""
    
    # Собираем ВСЕ вложения: и из текущего запроса, и из истории сообщений
    all_attachments = []
    if attachments:
        all_attachments.extend(attachments)
    
    # Ищем вложения в истории сообщений
    for m in messages:
        if m.get('attachments'):
            # Избегаем дубликатов, если вложения из текущего запроса уже попали в историю
            for att in m['attachments']:
                if not any(a.get('name') == att.get('name') for a in all_attachments):
                    all_attachments.append(att)

    if all_attachments:
        log.info(f"[DEBUG] Processing {len(all_attachments)} total attachments (current + history)")
        for att in all_attachments:
            try:
                att_name = att.get('name', 'file')
                # Данные могут быть в 'data' (base64) или в 'local_url' (путь на диске)
                att_data = att.get('data')
                local_url = att.get('local_url')
                
                file_bytes = None
                if att_data:
                    import base64
                    file_bytes = base64.b64decode(att_data)
                elif local_url:
                    # Если данных нет в base64, пробуем прочитать с диска
                    # Поддерживаем и старый путь (/img/) и новый (/uploads/)
                    if "/api/chat/uploads/" in local_url:
                        # /api/chat/uploads/{client_id}/chat_files/{session_id}/{filename}
                        parts = local_url.split("/api/chat/uploads/")[-1].split("/")
                        full_path = os.path.join(BASE_DIR, "uploads", *parts)
                    elif "/api/chat/img/" in local_url:
                        parts = local_url.split("/api/chat/img/")[-1].split("/")
                        full_path = os.path.join(BASE_DIR, "img", *parts)
                    else:
                        full_path = None

                    if full_path and os.path.exists(full_path):
                        with open(full_path, "rb") as f:
                            file_bytes = f.read()

                if file_bytes:
                    from io import BytesIO
                    ext = att_name.split('.')[-1].lower()
                    att_text = ""
                    if ext in ['txt', 'md', 'json', 'csv', 'yml', 'yaml', 'robots']:
                        att_text = file_bytes.decode('utf-8', errors='ignore')
                    elif ext == 'pdf':
                        import pdfplumber
                        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                            att_text = "\n".join([p.extract_text() or "" for p in pdf.pages])
                    elif ext == 'docx':
                        import docx
                        doc = docx.Document(BytesIO(file_bytes))
                        att_text = "\n".join([p.text for p in doc.paragraphs])
                    elif ext == 'xlsx':
                        import openpyxl
                        wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
                        rows = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                rows.append('\t'.join(str(c or '') for c in row))
                        att_text = '\n'.join(rows)
                    elif ext in ('ppt', 'pptx'):
                        from pptx import Presentation
                        prs = Presentation(BytesIO(file_bytes))
                        slides = []
                        for slide in prs.slides:
                            texts = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    texts.append(shape.text)
                            slides.append('\n'.join(texts))
                        att_text = '\n\n'.join(slides)
                    
                    if att_text:
                        user_attachments_content += f"\n--- ТЕКСТ ИЗ ФАЙЛА '{att_name}' ---\n{att_text}\n--- КОНЕЦ ФАЙЛА '{att_name}' ---\n"
            except Exception as e:
                log.error(f"Attachment Processing Error: {e}")

    # 2. СТУПЕНЬ: ЛИЧНОСТЬ И ПРОФИЛЬ (ДАННЫЕ ИЗ ПОЛЕЙ)
    contact_info = ""
    final_site_url = ""
    
    # Получаем контакты максимально надежно
    if isinstance(client_config, dict):
        contacts = client_config.get('contacts', {})
        final_site_url = client_config.get('site_url') or client_config.get('bot_settings', {}).get('site_url') or ""
    else:
        # Это объект ClientConfig
        contacts = client_config.raw.get('contacts', {})
        final_site_url = client_config.raw.get('site_url') or client_config.raw.get('bot_settings', {}).get('site_url') or ""

    # Очистка final_site_url от возможных артефактов
    if final_site_url:
        final_site_url = str(final_site_url).strip().rstrip(')')

    legal_data = client_config.get('legal_data', {}) if isinstance(client_config, dict) else client_config.raw.get('legal_data', {})

    # Затем добавляем основной файл знаний из настроек
    knowledge_file_url = bot_settings.get('knowledge_file_url')
    if knowledge_file_url and bot_settings.get('enable_knowledge_file', True):
        try:
            # Извлекаем путь из URL. Поддерживаем и старый формат (/img/) и новый (/uploads/)
            if "/uploads/" in knowledge_file_url:
                # Формат: /api/chat/uploads/{client_id}/knowledge/{filename}
                parts = knowledge_file_url.split("/uploads/")[-1].split("/")
                file_path = os.path.join(BASE_DIR, "uploads", *parts)
            else:
                # Старый формат: /api/chat/img/{client_id}/{filename} или /api/chat/img/{filename}
                filename = knowledge_file_url.split('/')[-1]
                # Пробуем сначала в подпапке клиента, потом в корне img
                file_path = os.path.join(BASE_DIR, 'img', client_id, filename)
                if not os.path.exists(file_path):
                    file_path = os.path.join(BASE_DIR, 'img', filename)

            if os.path.exists(file_path):
                ext = file_path.split('.')[-1].lower()
                raw_text = ""
                if ext in ['txt', 'md', 'json', 'csv', 'yml', 'yaml']:
                    with open(file_path, 'r', encoding='utf-8') as f: raw_text = f.read()
                elif ext == 'pdf':
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        raw_text = "\n".join([p.extract_text() for p in pdf.pages if p.extract_text()])
                elif ext == 'docx':
                    import docx
                    doc = docx.Document(file_path)
                    raw_text = "\n".join([p.text for p in doc.paragraphs])

                if raw_text:
                    # Используем векторный поиск
                    vector_db = VectorService(f"{client_id}:{assistant_id or 'main'}")
                    # Простая проверка: если в индексе нет данных, индексируем
                    if not vector_db.chunks:
                        await vector_db.add_texts([raw_text])
                    
                    # Ищем релевантные куски по последнему сообщению пользователя
                    user_query = messages[-1]['content'] if messages else ""
                    relevant_chunks = await vector_db.search(user_query, top_k=5)
                    
                    if relevant_chunks:
                        file_content = "\n### [ИНФОРМАЦИЯ ИЗ БАЗЫ ЗНАНИЙ (ФАЙЛЫ)]:\n" + "\n---\n".join(relevant_chunks) + "\n"
                        log.info(f"[Vector] Found {len(relevant_chunks)} relevant chunks for {client_id}")
                    else:
                        # Если ничего не нашли, но файл есть - берем начало файла как fallback
                        file_content = f"\n### [ГЛАВНЫЙ РЕГЛАМЕНТ (НАЧАЛО)]:\n{raw_text[:2000]}\n"
                
            else:
                log.warning(f"Knowledge file NOT FOUND at path: {file_path}")
        except Exception as e: log.error(f"File Error: {e}")

    if contacts:
        log.info(f"[DEBUG] Raw contacts from config: {contacts}")
        
        # Формируем блок контактов только если есть хоть какие-то данные
        temp_contact_info = "### [КОНТАКТНАЯ ИНФОРМАЦИЯ (ПРИОРИТЕТ №1)]:\n"
        has_any_contact = False

        if contacts.get('company_name'): 
            temp_contact_info += f"Компания: {contacts['company_name']}\n"
            has_any_contact = True
        
        # Собираем все телефоны
        phones_list = []
        main_phone = contacts.get('phone')
        if main_phone: 
            phones_list.append({"label": "Телефон", "value": main_phone, "link": f"tel:{re.sub(r'[^0-9+]', '', main_phone)}"})
        
        extra = contacts.get('extra_phones', [])
        if isinstance(extra, list):
            for item in extra:
                if isinstance(item, dict) and item.get('phone'):
                    p_val = item['phone']
                    phones_list.append({
                        "label": item.get('label', 'Телефон'), 
                        "value": p_val, 
                        "link": f"tel:{re.sub(r'[^0-9+]', '', p_val)}"
                    })
        
        if phones_list:
            temp_contact_info += "СПИСОК ТЕЛЕФОНОВ:\n"
            for p in phones_list: 
                temp_contact_info += f"Метка: {p['label']} | Номер: {p['value']} | Ссылка: {p['link']}\n"
            has_any_contact = True
        
        # Собираем все Email
        emails_list = []
        main_email = contacts.get('email')
        if main_email: 
            emails_list.append({"label": "Почта", "value": main_email, "link": f"mailto:{main_email}"})
        
        extra_emails = contacts.get('extra_emails', [])
        if isinstance(extra_emails, list) and extra_emails:
            for em in extra_emails:
                if isinstance(em, dict) and em.get('email'):
                    e_val = em['email']
                    emails_list.append({
                        "label": em.get('label', 'Почта'), 
                        "value": e_val, 
                        "link": f"mailto:{e_val}"
                    })
        
        if emails_list:
            temp_contact_info += "СПИСОК EMAIL:\n"
            for e in emails_list: 
                temp_contact_info += f"Метка: {e['label']} | Адрес: {e['value']} | Ссылка: {e['link']}\n"
            has_any_contact = True
        
        # Собираем все адреса
        addresses_list = []
        extra_addresses = contacts.get('extra_addresses', [])
        if isinstance(extra_addresses, list) and extra_addresses:
            for a in extra_addresses:
                if isinstance(a, dict) and a.get('address'):
                    addresses_list.append({
                        "label": a.get('label', 'Адрес'),
                        "value": a['address']
                    })
        
        if addresses_list:
            temp_contact_info += "СПИСОК АДРЕСОВ:\n"
            for a in addresses_list:
                temp_contact_info += f"Метка: {a['label']} | Адрес: {a['value']}\n"
            has_any_contact = True
        
        # Мессенджеры
        for m_name, m_key in [('WhatsApp', 'extra_wa'), ('Telegram', 'extra_tg'), ('Max', 'extra_max'), ('VK Messenger', 'extra_vk')]:
            items = contacts.get(m_key, [])
            if isinstance(items, list) and items:
                for it in items:
                    if isinstance(it, dict) and it.get('value'):
                        label = it.get('label', m_name)
                        val = it['value'].strip()
                        mode = it.get('mode', 'user')
                        
                        link = val
                        if m_name == 'WhatsApp' and mode == 'user':
                            phone_digits = re.sub(r'\D', '', val)
                            link = f"https://wa.me/{phone_digits}"
                        elif m_name == 'Telegram' and mode == 'user':
                            link = f"https://t.me/{val.replace('@', '')}"
                        elif m_name == 'VK Messenger' and mode == 'user':
                            link = f"https://vk.me/{val.replace('@', '')}"
                        elif not link.startswith('http'):
                            link = f"https://{link}"
                            
                        temp_contact_info += f"МЕССЕНДЖЕР: {m_name} | Метка: {label} | Тип: {mode} | Значение: {val} | Ссылка: {link}\n"
                        has_any_contact = True

        # Социальные сети (extra_links)
        extra_links = contacts.get('extra_links', [])
        if isinstance(extra_links, list) and extra_links:
            for link_item in extra_links:
                if isinstance(link_item, dict) and link_item.get('url'):
                    label = link_item.get('label', 'Соцсеть')
                    url = link_item['url'].strip()
                    if not url.startswith('http'):
                        url = f"https://{url}"
                    temp_contact_info += f"СОЦСЕТЬ: Метка: {label} | Ссылка: {url}\n"
                    has_any_contact = True

        if final_site_url: 
            temp_contact_info += f"Сайт: {final_site_url}\n"
            has_any_contact = True

        if has_any_contact:
            rules = [
                "\n### ПРАВИЛА ОФОРМЛЕНИЯ КОНТАКТОВ (СТРОГО):",
                "При ответе на вопрос о контактах, используй ТОЛЬКО этот формат:",
                "1. ТЕЛЕФОН: **Название**: [номер](ссылка). Пример: **Телефон**: [8 (999) 000-00-00](tel:+79990000000)",
                "2. ПОЧТА: **Почта**: [название почты](ссылка). Пример: **Почта**: [info@example.com](mailto:info@example.com)",
                "3. МЕССЕНДЖЕРЫ: **Название**: [username или номер](ссылка). Пример: **Telegram**: [@username](https://t.me/username)",
                "4. СОЦСЕТИ: **Название соцсети**: [название](ссылка). Пример: **VK Group**: [VK Group](https://vk.com/...)",
                "5. КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО использовать любые символы в начале строк (тире, точки, галочки, ✔️, •, -). Просто текст с новой строки.",
                "6. ЗАПРЕЩЕНО писать слово 'Ссылка' или 'Номер' перед самими данными.",
                "7. ЗАПРЕЩЕНО дублировать URL в тексте ссылки. Используй формат [название](url), где название — это метка из реестра, а не сам URL.",
                "8. Если в реестре указана конкретная метка (например, 'Отдел продаж'), используй её вместо стандартного 'Телефон'.",
                "9. ЗАПРЕЩЕНО выдумывать контакты, которых нет в списке выше! Если списка нет или он пуст — вежливо скажи, что контакты не указаны."
            ]
            contact_info = temp_contact_info + "\n".join(rules)



        
        # Определяем тип реквизитов (ИП, ООО, Самозанятый)
        legal_config = client_config.get('legal', {}) if is_dict else client_config.raw.get('legal', {})
        legal_type = legal_config.get('type', 'ip')
        
        type_data = legal_data.get(legal_type, {})
        if type_data:
            legal_labels = {
                'ip': {'name': 'Индивидуальный предприниматель', 'ogrn': 'ОГРНИП', 'entity': 'ИП'},
                'ooo': {'name': 'Юридическое лицо', 'ogrn': 'ОГРН', 'entity': 'Название'},
                'self': {'name': 'Самозанятый', 'ogrn': 'ИНН', 'entity': 'ФИО'}
            }
            label = legal_labels.get(legal_type, legal_labels['ip'])
            
            # Проверяем, есть ли хоть какие-то данные для вывода
            has_data = any(type_data.get(k) for k in ['name', 'inn', 'ogrn', 'address'])
            
            if has_data:
                contact_info += f"\n### [РЕКВИЗИТЫ {label['name'].upper()}]:\n"
                
                if legal_type == 'self':
                    contact_info += f"Статус: {label['name']}\n"
                    if type_data.get('name'): contact_info += f"ФИО: {type_data['name']}\n"
                    if type_data.get('inn'): contact_info += f"ИНН: {type_data['inn']}\n"
                    if type_data.get('birth_date'): contact_info += f"Дата рождения: {type_data['birth_date']}\n"
                    if type_data.get('birth_place'): contact_info += f"Место рождения: {type_data['birth_place']}\n"
                    if type_data.get('reg_address'): contact_info += f"Адрес регистрации: {type_data['reg_address']}\n"
                    
                    # ПАСПОРТНЫЕ ДАННЫЕ
                    if type_data.get('pass_seria') and type_data.get('pass_number'):
                        contact_info += f"Паспорт: серия {type_data['pass_seria']} № {type_data['pass_number']}\n"
                    if type_data.get('pass_date'):
                        contact_info += f"Дата выдачи паспорта: {type_data['pass_date']}\n"
                    if type_data.get('pass_code'):
                        contact_info += f"Код подразделения: {type_data['pass_code']}\n"
                    if type_data.get('pass_issuer'):
                        contact_info += f"Кем выдан паспорт: {type_data['pass_issuer']}\n"
                else:
                    contact_info += f"Статус: {label['name']}\n"
                    if type_data.get('name'): contact_info += f"Название: {type_data['name']}\n"
                    if type_data.get('ogrn'): contact_info += f"{label['ogrn']}: {type_data['ogrn']}\n"
                    if type_data.get('inn'): contact_info += f"ИНН: {type_data['inn']}\n"
                    if type_data.get('address'): contact_info += f"Юридический адрес: {type_data['address']}\n"
                    if type_data.get('bank_name'): contact_info += f"Банк: {type_data['bank_name']}\n"
                    if type_data.get('bank_bik'): contact_info += f"БИК: {type_data['bank_bik']}\n"
                    if type_data.get('bank_account'): contact_info += f"Расчётный счёт: {type_data['bank_account']}\n"
                    if type_data.get('bank_corr'): contact_info += f"Корреспондентский счёт: {type_data['bank_corr']}\n"
                
                contact_info += (
                    "\n### ПРАВИЛА ОФОРМЛЕНИЯ РЕКВИЗИТОВ (СТРОГО):\n"
                    "При выводе реквизитов используй формат: **Название поля**: значение\n"
                    "Пример: **ИНН**: 1234567890\n"
                    "Пример: **Расчётный счёт**: 40802810XXXXXXXXXXXX\n"
                )

    # 2.1. ГРАФИК РАБОТЫ
    working_hours_info = ""
    
    # ПРАЗДНИКИ И ИСКЛЮЧЕНИЯ
    holidays_enabled = client_config.get('working_hours_holidays_enabled') if is_dict else client_config.raw.get('working_hours_holidays_enabled')
    holidays = client_config.get('working_hours_holidays') if is_dict else client_config.raw.get('working_hours_holidays')
    
    working_hours = client_config.get('working_hours', {}) if is_dict else client_config.raw.get('working_hours', {})
    if working_hours or (holidays_enabled and holidays):
        working_hours_info += "### [ГРАФИК РАБОТЫ И ПРАЗДНИКИ (СТРОГО)]:\n"
        
        # Словарь для перевода дня недели на русский
        days_ru = {
            'Monday': 'Понедельник', 'Tuesday': 'Вторник', 'Wednesday': 'Среда',
            'Thursday': 'Четверг', 'Friday': 'Пятница', 'Saturday': 'Суббота', 'Sunday': 'Воскресенье'
        }
        curr_day_en = time.strftime('%A', time.localtime())
        curr_day_ru = days_ru.get(curr_day_en, curr_day_en)
        
        working_hours_info += f"Сегодняшняя дата: {time.strftime('%d.%m.%Y', time.localtime())} ({curr_day_ru})\n\n"
        
        if holidays_enabled and holidays:
            working_hours_info += "КРИТИЧЕСКИ ВАЖНО (ПРАЗДНИКИ И ПЕРЕНОСЫ):\n"
            working_hours_info += "ВНИМАНИЕ! Если сегодня праздник из списка ниже, обычный график ИГНОРИРУЕТСЯ:\n"
            working_hours_info += f"{holidays}\n"
            working_hours_info += "Проверь сегодняшнюю дату и список праздников выше. Если совпадает — отвечай по праздничному графику.\n\n"

        active_days = []
        full_days_map = {'mon': 'Понедельник', 'tue': 'Вторник', 'wed': 'Среда', 'thu': 'Четверг', 'fri': 'Пятница', 'sat': 'Суббота', 'sun': 'Воскресенье'}
        for day_id, day_name in full_days_map.items():
            day_data = working_hours.get(day_id)
            if day_data and day_data.get('enabled'):
                is_24h = day_data.get('is_24h') or (day_data.get('from') == "00:00" and day_data.get('to') in ["00:00", "23:59", "24:00"])
                time_str = "Круглосуточно" if is_24h else f"с {day_data.get('from')} до {day_data.get('to')}"
                if day_data.get('lunch_enabled'):
                    time_str += f" (перерыв на обед с {day_data.get('lunch_from')} до {day_data.get('lunch_to')})"
                active_days.append(f"**{day_name}**: {time_str}")
        
        if active_days:
            working_hours_info += "ОБЫЧНЫЙ РАБОЧИЙ ГРАФИК:\n" + "\n".join(active_days) + "\n"
            
            weekend_days = [day_name for day_id, day_name in full_days_map.items() if not working_hours.get(day_id, {}).get('enabled')]
            if weekend_days: 
                working_hours_info += "ВЫХОДНЫЕ ДНИ: " + ", ".join([f"**{d}**" for d in weekend_days]) + "\n"
        
        working_hours_info += (
            "\n### ПРАВИЛА ОФОРМЛЕНИЯ ГРАФИКА (СТРОГО):\n"
            "1. Используй формат: **День недели**: время работы\n"
            "2. Если есть ПРАЗДНИКИ И ИСКЛЮЧЕНИЯ, обязательно учитывай их при ответе на вопрос 'как вы работаете сегодня?'.\n"
            "3. Если данных о графике нет — не упоминай его.\n"
            "4. ЗАПРЕЩЕНО использовать любые эмодзи (📅, 🕒) в графике.\n"
        )

    # 3. СТУПЕНЬ: ПРОИНДЕКСИРОВАННЫЕ СТРАНИЦЫ САЙТА (ПРИОРИТЕТ №3)
    site_context = ""
    if bot_settings.get('enable_site_search', True):
        try:
            from .site_indexer import get_indexer_for_client
            # Основной URL хранится в карточке ассистента. Поля bot_settings и contacts
            # поддерживаются только для обратной совместимости.
            site_url = final_site_url or bot_settings.get('site_url') or contacts.get('website')
            user_query = next((m['content'] for m in reversed(messages) if m['role'] == 'user'), "")
            if site_url and user_query:
                indexer = get_indexer_for_client(client_id, site_url, assistant_id=assistant_id)
                results = await indexer.search(user_query, limit=3)
                if results:
                    site_context = "\n### [ПРОИНДЕКСИРОВАННЫЕ СТРАНИЦЫ САЙТА (ПРИОРИТЕТ №3)]:\n" + "\n".join(
                        f"Источник: {r['title']} ({r['url']})\nТекст: {r['snippet']}" for r in results
                    )
                    log.info(f"[SiteSearch] Found {len(results)} pages for {client_id}:{assistant_id or 'main'}")
        except Exception as e:
            log.error(f"RAG Error for {client_id}:{assistant_id or 'main'}: {e}")

    # ПАРАМЕТРЫ ДНК
    dna_rules = []
    
    # Личность и инструкции (только если заполнены)
    personality = bot_settings.get('personality_prompt')
    if personality: dna_rules.append(f"О КОМПАНИИ/ПРОЕКТЕ: {personality}")
    
    instructions = bot_settings.get('negative_prompt')
    if instructions: dna_rules.append(f"ГЛАВНЫЕ ИНСТРУКЦИИ И ПРАВИЛА: {instructions}")

    
    # 1. Кто ведет диалог (dna_proactive)
    proactive = bot_settings.get('dna_proactive', 'reactive')
    if proactive == 'active':
        dna_rules.append("Будь проактивным: старайся вести диалог, задавай уточняющие вопросы и предлагай помощь, не дожидаясь инициативы пользователя.")
    else:
        dna_rules.append("Будь реактивным: отвечай строго на вопросы пользователя, не навязывай лишнюю информацию.")

    # 2. Обращение (dna_addressing)
    addr = bot_settings.get('dna_addressing', 'formal')
    if addr == 'informal':
        dna_rules.append("ПРАВИЛО ОБРАЩЕНИЯ: Общайся с клиентом строго на 'ты'. Используй 'ты', 'тебе', 'твой'. КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО использовать 'Вы', 'Вам', 'Ваш'.")
    else:
        dna_rules.append("ПРАВИЛО ОБРАЩЕНИЯ: Общайся с клиентом строго на 'Вы'. Соблюдай деловую этику.")
    
    # 3. Общение (dna_tone)
    tone = bot_settings.get('dna_tone', 'strict')
    if tone == 'soft':
        dna_rules.append("ТОН ОБЩЕНИЯ: Приветливый, дружелюбный, теплый, неформальное. Используй вежливые обороты, проявляй эмпатию.")
    else:
        dna_rules.append("ТОН ОБЩЕНИЯ: Официальный, строго деловая этика, лаконичный. Никаких лишних эмоций.")
    
    # 4. Терминалогия (dna_language)
    lang_complexity = bot_settings.get('dna_language', 'simple')
    if lang_complexity == 'expert':
        dna_rules.append("Используй профессиональную терминологию и глубокие экспертные пояснения.")
    else:
        dna_rules.append("Объясняй всё максимально просто и понятно, избегай сложных терминов.")

    # 5. Длина ответов (dna_length)
    length = bot_settings.get('dna_length', 'short')
    dna_rules.append("Отвечай коротко и по делу." if length == 'short' else "Отвечай максимально подробно и развернуто.")

    # 6. Вариативность речи не должна разрешать выдумывать факты о бизнесе.
    try:
        temp_float = min(max(float(bot_settings.get('temperature', 0.3)), 0.0), 0.3)
    except (TypeError, ValueError):
        temp_float = 0.3
    dna_rules.append("Формулируй ответ естественно, но используй только предоставленные источники. Если прямого ответа нет — сообщи об этом и не выдумывай информацию.")

    # 7. Использование эмодзи (dna_emojis)
    # Если True или 'none' (старый дефолт) — разрешаем
    emojis_val = bot_settings.get('dna_emojis', True)
    is_enabled = emojis_val in [True, 'none', 'yes', 'on']
    
    if is_enabled:
        dna_rules.append("Используй эмодзи, чтобы сделать общение более живым и дружелюбным.")
    else:
        dna_rules.append("КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО использовать любые эмодзи. Только текст.")

    # 8. Главный приоритет (dna_focus)
    focus = bot_settings.get('dna_focus', 'facts')
    if focus == 'client':
        dna_rules.append("Твой главный приоритет — забота о клиенте и его комфорте, даже если это требует более мягкого отхода от регламента.")
    else:
        dna_rules.append("Твой главный приоритет — точность и достоверность данных. Ни в коем случае не искажай факты ради вежливости.")

    # ЛОГ ДЛЯ ПРОВЕРКИ ДНК
    log.info(f"[DNA DEBUG] Active rules for {client_id}: {dna_rules} | Temperature: {temp_float}")

    # СБОРКА СИСТЕМНОГО ПРОМПТА
    bot_name = bot_settings.get('bot_name') or (client_config.get('bot_name') if is_dict else client_config.raw.get('bot_name')) or 'Митя'
    bot_role = bot_settings.get('bot_role') or 'ИИ-ассистент'
    
    system_content = f"ТВОЕ ИМЯ: {bot_name}.\n"
    system_content += f"ТВОЯ РОЛЬ: {bot_role}.\n"
    system_content += f"ТЫ — {bot_role.upper()}. Никогда не называй себя GigaChat или ИИ-моделью.\n"
    system_content += "КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО переименовывать мессенджер 'Max' в 'Максим' или любые другие имена. Это название бренда.\n"
    system_content += f"Если тебя спрашивают, как тебя зовут, ты отвечаешь: '{bot_name}'.\n"
    system_content += "Твоя задача — представлять интересы компании.\n\n"
    
    system_content += "### [ПРАВИЛА ДОСТОВЕРНОСТИ]\n"
    system_content += "1. Используй только те факты о бизнесе, услугах и ценах, которые указаны в разделах ниже.\n"
    system_content += "2. Если в предоставленных данных нет прямого ответа — вежливо сообщи об этом: 'К сожалению, у меня нет информации по данному вопросу. Оставьте ваши контакты, и менеджер уточнит детали'.\n"
    system_content += "3. Не выдумывай номера телефонов, адреса или ссылки.\n"
    system_content += "4. Твой единственный источник правды — этот системный промпт. История диалога может содержать устаревшие данные.\n\n"
    
    system_content += "### [ВАЛИДАЦИЯ КОНТАКТОВ]\n"
    system_content += "1. Если пользователь оставляет почту, телефон или ссылку, обязательно проверяй их корректность.\n"
    system_content += "2. Почта должна содержать '@' и домен (например, user@example.com). Если пользователь написал '@domain.com' или 'user@domain' без точки — это ошибка.\n"
    system_content += "3. Телефон должен содержать минимум 10-11 цифр. Если цифр меньше — это ошибка.\n"
    system_content += "4. Если формат неверный, вежливо укажи на ошибку и попроси уточнить данные. Не делай вид, что всё в порядке.\n\n"
    
    system_content += "### [СТРОГИЙ ПОРЯДОК ИСТОЧНИКОВ]\n"
    system_content += f"Последнее обновление: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(client_config.get('updated_at', time.time()) if is_dict else client_config.raw.get('updated_at', time.time())))}.\n"
    system_content += "1. База знаний из загруженного файла — основной источник для услуг, тарифов, цен и условий.\n"
    system_content += "2. Заполненные карточки (контакты, соцсети, почта, телефоны, адреса, график) — единственный источник этих реквизитов.\n"
    system_content += "3. Проиндексированные страницы сайта — используй только после первых двух источников и только как дополнение.\n"
    system_content += "4. Если источники не содержат точного факта, цены, условия или контакта — прямо скажи, что данных нет, и предложи связаться с менеджером. Никогда не предполагай и не создавай данные.\n"
    system_content += "5. Если данные в истории диалога противоречат этому промпту — игнорируй историю и отвечай по новым данным.\n\n"
    
    # Собираем в порядке приоритета
    if file_content: system_content += file_content
    if contact_info: system_content += contact_info
    if working_hours_info: system_content += working_hours_info
    if site_context: system_content += site_context
    
    # КОНТЕКСТ ОБЪЯВЛЕНИЯ (Avito, Юла и др. площадки)
    context = kwargs.get("context") or {}
    
    # Если context - это объект Pydantic (например, ChatContext), преобразуем в dict
    if hasattr(context, "dict"):
        context = context.dict()
    elif hasattr(context, "model_dump"):
        context = context.model_dump()
    
    item_title = context.get("title") if isinstance(context, dict) else getattr(context, "title", "")
    item_price = context.get("price") if isinstance(context, dict) else getattr(context, "price", "")
    item_url = context.get("url") if isinstance(context, dict) else getattr(context, "url", "")
    if item_title or item_price or item_url:
        item_block = "\n### [ИНФОРМАЦИЯ ОБ ОБЪЯВЛЕНИИ/ВАКАНСИИ (КОНТЕКСТ ОБРАЩЕНИЯ)]:\n"
        if item_title:
            item_block += f"Название: {item_title}\n"
        if item_price:
            item_block += f"Цена: {item_price}\n"
        if item_url:
            item_block += f"Ссылка: {item_url}\n"
        item_block += "\nПользователь обращается к вам по этому объявлению/вакансии. Отвечай с учётом этого контекста.\n"
        system_content += item_block
    
    # Файлы пользователя добавляем в самый конец, чтобы они были максимально близки к вопросу
    if user_attachments_content: 
        system_content += f"\n### [АКТУАЛЬНЫЕ ФАЙЛЫ ОТ ПОЛЬЗОВАТЕЛЯ]:\n{user_attachments_content}\n"
        system_content += "ИНСТРУКЦИЯ: Пользователь прислал тебе файлы выше. Если он спрашивает о них, используй этот текст для ответа.\n"

    
    system_content += "\n### [ПРАВИЛА ОБЩЕНИЯ]:\n- " + "\n- ".join(dna_rules) + "\n"
    
    system_content += "\n### [ИНСТРУКЦИЯ ПО ОТВЕТАМ]:\n"
    system_content += "- Если пользователь спрашивает 'КОНТАКТЫ', предоставь доступные способы связи (телефон, почта, мессенджеры). Не выводи реквизиты или график работы, если о них не спросили отдельно.\n"
    system_content += "- Если пользователь спрашивает 'РЕКВИЗИТЫ', выводи полные данные из секции РЕКВИЗИТЫ выше.\n"
    system_content += "- Если пользователь спрашивает 'РЕЖИМ РАБОТЫ' или 'ГРАФИК', отвечай строго по секции РЕЖИМ РАБОТЫ выше.\n"
    system_content += "- ТЕБЕ РАЗРЕШЕНО поддерживать легкую беседу (приветствия, вопросы о погоде, настроении, кто ты такой). Отвечай вежливо и кратко, используя свои внутренние знания, но после ответа ОБЯЗАТЕЛЬНО делай мягкий переход к услугам компании.\n"

    # Приветствие — только в самом начале диалога. Если ассистент уже отвечал в этой
    # переписке, повторно здороваться нельзя (частая жалоба: бот "здоровается каждый раз").
    already_greeted = any(m.get('role') == 'assistant' for m in messages)
    if already_greeted:
        system_content += "- КРИТИЧЕСКИ ВАЖНО: диалог уже идёт, ты здоровался ранее. НЕ здоровайся повторно ('Здравствуйте', 'Добрый день', 'Привет' и т.п.) и не представляйся заново. Сразу отвечай по существу вопроса.\n"
    else:
        system_content += "- Это начало диалога: поздоровайся один раз, коротко и естественно, без формальных штампов.\n"
    
    system_content += "\nКРИТИЧЕСКОЕ ПРАВИЛО: Твои знания о БИЗНЕСЕ ограничены СТРОГО предоставленным РЕГЛАМЕНТОМ, ФАЙЛАМИ, КОНТАКТАМИ, ГРАФИКОМ РАБОТЫ и САЙТОМ. Если в них нет прямого ответа — КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО выдумывать факты. Вместо этого вежливо сообщи, что информация не указана. Всегда отдавай приоритет информации из РЕГЛАМЕНТА (Приоритет №1). Если там нет ответа, смотри ПРОФИЛЬ (Приоритет №2), затем ГРАФИК РАБОТЫ, затем САЙТ (Приоритет №3)."

    # Логируем только размер промпта (без содержимого — там могут быть ПДн и реквизиты)
    log.info(f"System prompt built for {client_id}: {len(system_content)} chars")

    has_system = False
    for m in messages:
        if m.get('role') == 'system':
            m['content'] = system_content
            has_system = True
            break
    
    if not has_system:
        messages.insert(0, {"role": "system", "content": system_content})

    # Получаем время последнего обновления конфига
    last_update = client_config.get('updated_at', 0) if is_dict else client_config.raw.get('updated_at', 0)

    if len(messages) >= 2 and messages[-1]['role'] == 'user':
        current_content = messages[-1]['content']
        reminder_tag = "[АКТУАЛЬНАЯ УСТАНОВКА:"
        if reminder_tag not in current_content:
            emoji_reminder = "Пиши ТОЛЬКО ТЕКСТОМ, КАТЕГОРИЧЕСКИ ЗАПРЕЩЕНО использовать любые эмодзи и иконки (📅, 🕒, ✅ и др.)." if not is_enabled else "Можешь использовать эмодзи для живости общения."
            
            # Собираем только те данные, которые явно указаны в панели управления
            overrides = []
            if final_site_url: 
                overrides.append(f"Сайт: {final_site_url}")
            
            # Собираем телефоны из динамического списка
            extra_p = contacts.get('extra_phones', [])
            if isinstance(extra_p, list) and extra_p:
                phones = [f"{p['label']}: {p['phone']}" for p in extra_p if isinstance(p, dict) and p.get('phone')]
                if phones: overrides.append("Телефоны: " + ", ".join(phones))
            
            # Собираем почты из динамического списка
            extra_e = contacts.get('extra_emails', [])
            if isinstance(extra_e, list) and extra_e:
                emails = [f"{e['label']}: {e['email']}" for e in extra_e if isinstance(e, dict) and e.get('email')]
                if emails: overrides.append("Email: " + ", ".join(emails))

            # Собираем адреса из динамического списка
            extra_a = contacts.get('extra_addresses', [])
            if isinstance(extra_a, list) and extra_a:
                addresses = [f"{a['label']}: {a['address']}" for a in extra_a if isinstance(a, dict) and a.get('address')]
                if addresses: overrides.append("Адреса: " + ", ".join(addresses))

            # Собираем мессенджеры для напоминания
            for m_name, m_key in [('WhatsApp', 'extra_wa'), ('Telegram', 'extra_tg'), ('Max', 'extra_max'), ('VK', 'extra_vk')]:
                items = contacts.get(m_key, [])
                if isinstance(items, list) and items:
                    m_vals = [f"{it.get('label', m_name)}: {it['value']}" for it in items if isinstance(it, dict) and it.get('value')]
                    if m_vals: overrides.append(f"{m_name}: " + ", ".join(m_vals))

            # Собираем соцсети для напоминания
            extra_links = contacts.get('extra_links', [])
            if isinstance(extra_links, list) and extra_links:
                links = [f"{l['label']}: {l['url']}" for l in extra_links if isinstance(l, dict) and l.get('url')]
                if links: overrides.append("Соцсети: " + ", ".join(links))
            
            # Добавляем краткую сводку реквизитов
            legal_config = client_config.get('legal', {}) if is_dict else client_config.raw.get('legal', {})
            l_type = legal_config.get('type', 'ip')
            l_data = legal_data.get(l_type, {})
            if l_data.get('name'):
                overrides.append(f"Реквизиты: {l_data.get('name')}")
            
            overrides_str = ". ".join(overrides) + "." if overrides else "Контактные данные не указаны."
            
            # Метка времени для борьбы с контекстом
            update_time = time.strftime('%H:%M:%S', time.localtime(last_update)) if last_update > 0 else ""
            time_str = f" (данные актуальны на {update_time})" if update_time else ""

            reminder = f"\n\n{reminder_tag} Твое имя {bot_name}. Твоя роль: {bot_role}. {overrides_str}{time_str} {emoji_reminder} Если контакта нет в списке выше — его не существует. Не выдумывай данные. Используй актуальную информацию из системного промпта.]"
            messages[-1]['content'] += reminder


    # 6. ВЫЗОВ МОДЕЛИ
    vision_file_ids = []
    # Для vision передаём только файлы текущего сообщения, не загружаем историю повторно.
    for attachment in attachments or []:
        content_type = str(attachment.get('content_type') or attachment.get('type') or '').lower()
        attachment_name = str(attachment.get('name') or attachment.get('file_name') or '')
        is_image = content_type.startswith('image/') or attachment_name.lower().endswith(('.jpg', '.jpeg', '.png', '.tiff', '.bmp'))
        local_url = attachment.get('local_url')
        if not is_image or not local_url or '/api/chat/uploads/' not in local_url:
            continue
        relative_path = local_url.split('/api/chat/uploads/', 1)[1]
        image_path = os.path.join(BASE_DIR, 'uploads', *relative_path.split('/'))
        if not os.path.isfile(image_path) or os.path.getsize(image_path) > 15 * 1024 * 1024:
            continue
        try:
            from .gigachat_service import upload_gigachat_file
            file_id = await upload_gigachat_file(
                image_path,
                attachment_name or os.path.basename(image_path),
                content_type or 'image/jpeg',
            )
            if file_id:
                vision_file_ids.append(file_id)
                break
        except Exception as error:
            log.warning('Vision upload skipped for %s: %s', attachment_name, error)

    if vision_file_ids and messages:
        messages[-1]['attachments'] = vision_file_ids
        messages[-1]['content'] += '\n\n[Системная заметка: проанализируй прикреплённое изображение и отвечай только по тому, что на нём действительно видно.]'

    rules = kwargs.get('rules', {})
    
    # Определяем модель на основе настроек бота и ограничений тарифа
    selected_model = bot_settings.get('ai_model', 'gigachat').lower()
    available_models = rules.get('available_models', ['GigaChat'])
    default_model = rules.get('model', 'GigaChat')

    # Анализ изображений доступен через Files API только в мультимодальной GigaChat-Pro.
    if vision_file_ids:
        target_model = 'GigaChat-Pro'
    elif 'yandex' in selected_model:
        target_model = 'yandexgpt/latest'
    else:
        target_model = 'GigaChat'

    # Если выбранная модель не разрешена тарифом, откатываемся на дефолтную для тарифа
    if target_model not in available_models and not vision_file_ids:
        log.warning(f"Model {target_model} not allowed for tariff. Falling back to {default_model}")
        target_model = default_model

    # Добавляем системные инструменты (время и др.)
    kwargs['available_tools'] = get_bot_tools(bot_settings)
    
    # Добавляем температуру в kwargs для передачи в сервисы
    kwargs['temperature'] = temp_float
    
    # ПОЛУЧАЕМ ЧЕРНОВИК ОТВЕТА
    if 'yandex' in target_model.lower():
        draft_response = await handle_yandex_request(messages, stream=stream, target_model=target_model, **kwargs)
    else:
        draft_response = await handle_gigachat_request(messages, target_model, stream=stream, **kwargs)

    # 7. DATA GUARD: собираем белый список контактов ВСЕГДА (и для стрима тоже)
    c_data = contacts if isinstance(contacts, dict) else {}
    allowed_list = []
    
    # Добавляем реквизиты в белый список, чтобы их не вырезал Data Guard
    l_config = client_config.get('legal', {}) if is_dict else client_config.raw.get('legal', {})
    l_type = l_config.get('type', 'ip')
    l_data = legal_data.get(l_type, {})
    for key in ['inn', 'ogrn', 'bank_bik', 'bank_account', 'bank_corr']:
        val = l_data.get(key)
        if val: 
            # Добавляем и чистые цифры, и с пробелами
            clean_val = str(val).replace(' ', '')
            allowed_list.append(clean_val)
            if ' ' in str(val):
                allowed_list.append(str(val))

    # Основной телефон
    main_phone = c_data.get('phone', '')
    if main_phone: allowed_list.append(main_phone)
    # Доп. телефоны
    for p in c_data.get('extra_phones', []):
        if isinstance(p, dict) and p.get('phone'): allowed_list.append(p['phone'])
    
    # Основная почта
    main_email = c_data.get('email', '')
    if main_email: allowed_list.append(main_email)
    # Доп. почты
    for e in c_data.get('extra_emails', []):
        if isinstance(e, dict) and e.get('email'): allowed_list.append(e['email'])

    # Адреса
    for a in c_data.get('extra_addresses', []):
        if isinstance(a, dict) and a.get('address'): allowed_list.append(a['address'])
        
    # Telegram (поле value, не username!)
    for t in c_data.get('extra_tg', []):
        if isinstance(t, dict):
            val = t.get('value') or t.get('username') or ''
            if val: allowed_list.append(val)

    # WhatsApp
    for w in c_data.get('extra_wa', []):
        if isinstance(w, dict):
            val = w.get('value') or ''
            if val: allowed_list.append(val)

    # Max
    for m in c_data.get('extra_max', []):
        if isinstance(m, dict):
            val = m.get('value') or ''
            if val: allowed_list.append(val)

    # VK
    for v in c_data.get('extra_vk', []):
        if isinstance(v, dict):
            val = v.get('value') or ''
            if val: allowed_list.append(val)

    # Соцсети
    for l in c_data.get('extra_links', []):
        if isinstance(l, dict) and l.get('url'): allowed_list.append(l['url'])
        
    if final_site_url: 
        allowed_list.append(final_site_url)

    log.info(f"[DATA GUARD] Allowed contacts for {client_id}: {allowed_list}")

    # Прокидываем allowed_list в kwargs для стрим-обработчиков
    kwargs['allowed_contacts'] = allowed_list

    # 8. DATA GUARD И АГЕНТ-КОНТРОЛЕР (только для не-стрима)
    if not stream and draft_response and isinstance(draft_response, str) and len(draft_response) > 5:
        # ПРИМЕНЯЕМ DATA GUARD
        draft_response = apply_data_guard(draft_response, allowed_list)

        # Собираем "Золотой реестр" для агента
        registry_parts = []
        if file_content: registry_parts.append(file_content)
        if contact_info: registry_parts.append(contact_info)
        if working_hours_info: registry_parts.append(working_hours_info)
        if site_context: registry_parts.append(site_context)
        if personality: registry_parts.append(f"О КОМПАНИИ: {personality}")
        if instructions: registry_parts.append(f"ПРАВИЛА: {instructions}")
        # Контекст объявления для агента
        context = kwargs.get("context") or {}
        item_title = context.get("title") or ""
        item_price = context.get("price") or ""
        item_url = context.get("url") or ""
        if item_title or item_price or item_url:
            item_block = "### [ИНФОРМАЦИЯ ОБ ОБЪЯВЛЕНИИ/ВАКАНСИИ (КОНТЕКСТ ОБРАЩЕНИЯ)]:\n"
            if item_title:
                item_block += f"Название: {item_title}\n"
            if item_price:
                item_block += f"Цена: {item_price}\n"
            if item_url:
                item_block += f"Ссылка: {item_url}\n"
            registry_parts.append(item_block)
        
        registry_data = "\n\n".join(registry_parts)
        
        if registry_data:
            log.info(f"[AGENT] Starting fact-check for {client_id}...")
            final_response = await validate_response_with_agent(draft_response, registry_data, **kwargs)
            # После агента ЕЩЕ РАЗ прогоняем через Data Guard для 100% гарантии
            final_response = apply_data_guard(final_response, allowed_list)
            log.info(f"[AGENT] Fact-check complete.")
            return final_response

    return draft_response

async def generate_faq(client_id: str, history_text: str):
    """Генерирует FAQ на основе истории диалогов.

    Returns:
        tuple: (frequent_requests, traffic_quality, spam_detected)
            - frequent_requests: list of dicts [{"question": "...", "count": N}, ...]
            - traffic_quality: str с анализом качества трафика
            - spam_detected: list of dicts с обнаруженным спамом
    """
    token = await get_gigachat_token()
    if not token:
        return [], "Не удалось авторизоваться в GigaChat", []

    prompt = f"""Ты — аналитик поддержки. Проанализируй историю диалогов с клиентами и выдай СТРОГО JSON:

{{
  "frequent_requests": [
    {{"question": "краткая формулировка вопроса", "count": число_повторений}}
  ],
  "traffic_quality": "краткий анализ качества трафика (1-2 предложения)",
  "spam_detected": [
    {{"pattern": "описание спам-паттерна", "count": число_случаев}}
  ]
}}

Правила:
- frequent_requests: не более 10 самых частых вопросов. count — примерное число.
- traffic_quality: бизнес-ли это трафик или случайные посетители.
- spam_detected: если спама нет, верни пустой список [].
- вопросы пиши на русском, в точности как спрашивали клиенты.

ИСТОРИЯ ДИАЛОГОВ (только сообщения клиентов):
{history_text[:8000]}
"""

    try:
        raw_response = await handle_gigachat_request(
            [
                {"role": "system", "content": "Ты — аналитик. Отвечай строго в формате JSON без лишнего текста."},
                {"role": "user", "content": prompt}
            ],
            GIGACHAT_MODEL,
            stream=False,
            temperature=0.3
        )

        json_match = re.search(r'\{.*\}', raw_response, re.DOTALL)
        if json_match:
            data = json.loads(json_match.group(0))
            frequent_requests = data.get('frequent_requests', [])
            traffic_quality = data.get('traffic_quality', 'Недостаточно данных для анализа.')
            spam_detected = data.get('spam_detected', [])
            return frequent_requests, traffic_quality, spam_detected
        else:
            log.warning(f"generate_faq: could not parse JSON from response: {raw_response[:200]}")
            return [], "Не удалось распознать ответ ИИ.", []
    except Exception as e:
        log.error(f"generate_faq error: {e}")
        return [], f"Ошибка при генерации FAQ: {e}", []
